from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new Presentation
prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height


def add_slide(title_text, content_text, bg_image_path):
    """
    Creates a new slide with a background image and a centered text-box for title and content.
    """
    # Use a blank slide layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Insert background image if available
    try:
        slide.shapes.add_picture(bg_image_path, 0, 0, width=slide_width, height=slide_height)
    except Exception as e:
        print(f"Could not load background image {bg_image_path}: {e}")

    # Create a textbox for slide content
    # Adjust left, top, width, height to control placement.
    left = Inches(1)
    top = Inches(1)
    width = slide_width - Inches(2)
    height = slide_height - Inches(2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    # Title Paragraph
    title_para = text_frame.add_paragraph()
    title_para.text = title_text
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Spacer Paragraph
    spacer = text_frame.add_paragraph()
    spacer.text = ""

    # Content Paragraph
    content_para = text_frame.add_paragraph()
    content_para.text = content_text.strip()
    content_para.font.size = Pt(28)
    content_para.alignment = PP_ALIGN.CENTER


# Define background image paths (update these paths to match your files)
bg_images = [
    "background1.jpg",  # for Starter Task
    "background2.jpg",  # for Vocabulary
    "background3.jpg",  # for Retention and Recall
    "background4.jpg",  # for Lesson Objective
    "background5.jpg",  # for Tasks
    "background6.jpg",  # for Stretch and Challenge
    "background7.jpg",  # for Plenary
    "background8.jpg"  # for Homework
]

# Slide 1: Starter Task
add_slide(
    "Starter Task",
    "Compare posters for children and adults.",
    bg_images[0]
)

# Slide 2: Vocabulary
add_slide(
    "Vocabulary",
    "Contrast, Legibility, Alignment, Composition",
    bg_images[1]
)

# Slide 3: Retention and Recall
add_slide(
    "Retention and Recall",
    "Review last week's design principles.",
    bg_images[2]
)

# Slide 4: Lesson Objective
add_slide(
    "Lesson Objective",
    "Learn how design differs based on the audience.",
    bg_images[3]
)

# Slide 5: Tasks
tasks_text = (
    "• Identify features of an effective children’s flyer.\n"
    "• Begin designing the flyer layout.\n"
    "• Select fonts and colours suitable for children."
)
add_slide(
    "Tasks",
    tasks_text,
    bg_images[4]
)

# Slide 6: Stretch and Challenge
add_slide(
    "Stretch and Challenge",
    "Research psychology of colours in design.",
    bg_images[5]
)

# Slide 7: Plenary
add_slide(
    "Plenary",
    "Share initial flyer drafts.",
    bg_images[6]
)

# Slide 8: Homework
add_slide(
    "Homework",
    "Analyse two real-world posters and explain their effectiveness.",
    bg_images[7]
)

# Save the presentation
pptx_filename = "Canva_Lesson_2_Understanding_Target_Audiences.pptx"
prs.save(pptx_filename)
print(f"Presentation saved as {pptx_filename}")
